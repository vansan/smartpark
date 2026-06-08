from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# UAE Flag Colors
UAE_GREEN  = RGBColor(0x00, 0xB0, 0x5C)
UAE_RED    = RGBColor(0xE0, 0x2A, 0x3C)
UAE_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
UAE_BLACK  = RGBColor(0x07, 0x0B, 0x08)
UAE_DARK   = RGBColor(0x10, 0x17, 0x12)
UAE_CARD   = RGBColor(0x16, 0x20, 0x19)
GREY_TEXT  = RGBColor(0x82, 0x90, 0x87)
LIGHT_GREY = RGBColor(0xF0, 0xF4, 0xF1)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]  # completely blank

# ── helpers ──────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, line=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=UAE_WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb

def uae_flag_bar(slide, y=0, h=0.08):
    """Thin UAE flag stripe at top: red | green | white | black"""
    w = prs.slide_width.inches
    seg = w / 4
    add_rect(slide, 0,       y, seg,   h, fill=UAE_RED)
    add_rect(slide, seg,     y, seg,   h, fill=UAE_GREEN)
    add_rect(slide, seg*2,   y, seg,   h, fill=UAE_WHITE)
    add_rect(slide, seg*3,   y, seg,   h, fill=UAE_BLACK)

def dark_bg(slide):
    add_rect(slide, 0, 0, 13.33, 7.5, fill=UAE_BLACK)

def slide_number(slide, n, total=11):
    add_text(slide, f"{n} / {total}", 12.2, 7.1, 1, 0.35,
             size=9, color=GREY_TEXT, align=PP_ALIGN.RIGHT)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
dark_bg(sl)

# Left green accent bar
add_rect(sl, 0, 0, 0.25, 7.5, fill=UAE_GREEN)

# Red accent bottom strip
add_rect(sl, 0, 6.8, 13.33, 0.7, fill=UAE_RED)

# Flag badge top-right
add_rect(sl, 11.5, 0.3, 0.08, 0.8, fill=UAE_RED)
add_rect(sl, 11.58, 0.3,  0.58, 0.27, fill=UAE_GREEN)
add_rect(sl, 11.58, 0.57, 0.58, 0.27, fill=UAE_WHITE)
add_rect(sl, 11.58, 0.84, 0.58, 0.27, fill=UAE_BLACK)

# Title
add_text(sl, "🅿  SmartPark", 0.6, 1.6, 10, 1.4, size=54, bold=True, color=UAE_WHITE)
add_text(sl, "Smart Parking Management System", 0.6, 3.0, 10, 0.7,
         size=24, color=UAE_GREEN)
add_text(sl, "Proof of Concept  ·  UAE  ·  2025", 0.6, 3.75, 10, 0.5,
         size=16, color=GREY_TEXT)
add_text(sl, "🌐  smartpark-lac.vercel.app", 0.6, 6.85, 8, 0.5,
         size=14, color=UAE_WHITE)
slide_number(sl, 1)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
dark_bg(sl)
uae_flag_bar(sl)

add_text(sl, "The Problem", 0.5, 0.25, 12, 0.7, size=32, bold=True, color=UAE_WHITE)
add_text(sl, "Parking today is broken", 0.5, 0.95, 12, 0.4, size=16, color=UAE_RED)

problems = [
    ("🚗", "Drivers circle 10-15 min",     "No visibility of free slots"),
    ("📄", "Paper tickets & manual entry",  "Human error, revenue leakage"),
    ("📵", "No digital booking system",     "Slots double-booked or wasted"),
    ("👮", "Guards check passes manually",  "Slow, error-prone, no audit trail"),
    ("⏱️", "Expired bookings stay locked",  "Good slots wasted, revenue lost"),
]
for i, (icon, prob, desc) in enumerate(problems):
    y = 1.5 + i * 1.05
    add_rect(sl, 0.5, y, 12.3, 0.85, fill=UAE_CARD,
             line=UAE_RED, line_width=Pt(1))
    add_text(sl, icon,  0.7,  y+0.05, 0.6,  0.7, size=22)
    add_text(sl, prob,  1.35, y+0.04, 5.5,  0.4, size=15, bold=True, color=UAE_WHITE)
    add_text(sl, desc,  1.35, y+0.44, 5.5,  0.35, size=12, color=GREY_TEXT)
    add_text(sl, "❌",  7.5,  y+0.15, 0.5,  0.5, size=18, color=UAE_RED)

slide_number(sl, 2)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — The Solution
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
dark_bg(sl)
uae_flag_bar(sl)

add_text(sl, "SmartPark Solution", 0.5, 0.25, 12, 0.7, size=32, bold=True, color=UAE_WHITE)
add_text(sl, "One platform · Three roles · Real-time", 0.5, 0.95, 12, 0.4,
         size=16, color=UAE_GREEN)

cols = [
    ("🟢", "Driver",  UAE_GREEN, ["Find free slots on live map",
                                    "Book by date, time & duration",
                                    "Prepaid · QR code entry pass",
                                    "Countdown timer on booking"]),
    ("🔴", "Operator", UAE_RED,  ["Real-time slot dashboard",
                                    "Full booking audit trail",
                                    "Auto-release expired slots",
                                    "VAT-compliant billing"]),
    ("⬜", "Guard",   UAE_WHITE, ["QR scanner on any device",
                                    "Instant booking validation",
                                    "Slot marked occupied on scan",
                                    "No hardware required"]),
]
for i, (icon, role, color, points) in enumerate(cols):
    x = 0.5 + i * 4.3
    add_rect(sl, x, 1.5, 4.0, 5.5, fill=UAE_CARD, line=color, line_width=Pt(1.5))
    add_text(sl, icon, x+0.2, 1.6, 0.6, 0.6, size=28)
    add_text(sl, role, x+0.85, 1.65, 3.0, 0.5, size=20, bold=True, color=color)
    for j, pt in enumerate(points):
        add_text(sl, f"✓  {pt}", x+0.25, 2.45 + j*0.9, 3.6, 0.75,
                 size=13, color=UAE_WHITE)

slide_number(sl, 3)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Live Demo Flow
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
dark_bg(sl)
uae_flag_bar(sl)

add_text(sl, "Complete Booking Flow", 0.5, 0.25, 12, 0.7, size=32, bold=True, color=UAE_WHITE)
add_text(sl, "End-to-end · Live demo", 0.5, 0.95, 12, 0.35, size=16, color=UAE_GREEN)

steps = [
    ("1", "Login",          "Google Sign-In",              UAE_GREEN),
    ("2", "Map",            "See Dubai parking locations",  UAE_GREEN),
    ("3", "Slot Grid",      "Pick a free slot (🟢)",        UAE_GREEN),
    ("4", "Date & Time",    "Select date, time, duration",  UAE_GREEN),
    ("5", "Payment",        "AED 5/hr + 5% VAT",           UAE_GREEN),
    ("6", "QR Code",        "Booking confirmed + QR pass",  UAE_GREEN),
    ("7", "Guard Scan",     "QR scanned → slot occupied",   UAE_GREEN),
]

for i, (num, title, desc, color) in enumerate(steps):
    col = i % 4
    row = i // 4
    x = 0.4 + col * 3.25
    y = 1.5 + row * 2.8
    add_rect(sl, x, y, 3.0, 2.3, fill=UAE_CARD, line=UAE_GREEN, line_width=Pt(1))
    add_rect(sl, x, y, 3.0, 0.55, fill=UAE_GREEN)
    add_text(sl, f"Step {num}", x+0.15, y+0.07, 2.7, 0.4,
             size=13, bold=True, color=UAE_BLACK)
    add_text(sl, title, x+0.15, y+0.65, 2.7, 0.5, size=15, bold=True, color=UAE_WHITE)
    add_text(sl, desc,  x+0.15, y+1.2,  2.7, 0.9, size=12, color=GREY_TEXT)
    if i < len(steps) - 1 and col < 3:
        add_text(sl, "→", x+3.05, y+0.85, 0.3, 0.5, size=18, color=UAE_GREEN,
                 align=PP_ALIGN.CENTER)

slide_number(sl, 4)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Payment & Pricing
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
dark_bg(sl)
uae_flag_bar(sl)

add_text(sl, "Prepaid Payment Model", 0.5, 0.25, 12, 0.7, size=32, bold=True, color=UAE_WHITE)
add_text(sl, "Transparent · VAT-compliant · Fraud-proof", 0.5, 0.95, 12, 0.35,
         size=16, color=UAE_GREEN)

# Left: pricing example
add_rect(sl, 0.5, 1.5, 5.8, 5.5, fill=UAE_CARD, line=UAE_GREEN, line_width=Pt(1.5))
add_text(sl, "Sample Invoice", 0.8, 1.65, 5.2, 0.5, size=18, bold=True, color=UAE_GREEN)
items = [
    ("Dubai Mall · Slot P03", ""),
    ("Date", "07 Jun 2025"),
    ("Start Time", "10:00 AM"),
    ("End Time", "12:00 PM"),
    ("Duration", "2 hours"),
    ("", ""),
    ("Parking (AED 5 × 2 hrs)", "AED 10.00"),
    ("VAT (5%)", "AED 0.50"),
    ("─────────────────", "──────────"),
    ("TOTAL", "AED 10.50"),
]
for j, (lbl, val) in enumerate(items):
    y = 2.3 + j * 0.42
    col = UAE_GREEN if lbl == "TOTAL" else (UAE_WHITE if lbl.startswith("──") else GREY_TEXT)
    sz  = 15 if lbl == "TOTAL" else 12
    bd  = lbl == "TOTAL"
    add_text(sl, lbl, 0.8,  y, 3.5, 0.38, size=sz, bold=bd, color=col)
    add_text(sl, val, 4.5,  y, 1.6, 0.38, size=sz, bold=bd,
             color=UAE_GREEN if lbl == "TOTAL" else UAE_WHITE, align=PP_ALIGN.RIGHT)

# Right: key points
add_rect(sl, 6.9, 1.5, 5.9, 5.5, fill=UAE_CARD, line=UAE_RED, line_width=Pt(1.5))
add_text(sl, "Why Prepaid?", 7.1, 1.65, 5.5, 0.5, size=18, bold=True, color=UAE_RED)
bullets = [
    ("💰", "No revenue leakage",        "Slot guaranteed only after payment"),
    ("🔒", "Fraud prevention",          "QR ties payment to specific slot & time"),
    ("🧾", "VAT compliance built-in",   "5% UAE VAT calculated automatically"),
    ("💳", "Gateway-ready",             "Plug in Stripe, Telr, or any UAE PSP"),
    ("📊", "Full audit trail",          "Every transaction logged with timestamp"),
]
for j, (icon, title, sub) in enumerate(bullets):
    y = 2.3 + j * 0.88
    add_text(sl, icon,  7.1,  y,      0.5, 0.5, size=20)
    add_text(sl, title, 7.65, y,      4.9, 0.35, size=14, bold=True, color=UAE_WHITE)
    add_text(sl, sub,   7.65, y+0.38, 4.9, 0.35, size=11, color=GREY_TEXT)

slide_number(sl, 5)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — QR & Guard Flow
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
dark_bg(sl)
uae_flag_bar(sl)

add_text(sl, "QR Entry & Guard Validation", 0.5, 0.25, 12, 0.7, size=32, bold=True, color=UAE_WHITE)
add_text(sl, "Zero paperwork · Instant validation · Any device", 0.5, 0.95, 12, 0.35,
         size=16, color=UAE_GREEN)

flow_steps = [
    ("📱", "Driver Books",     "Selects slot,\npays online"),
    ("✅", "QR Generated",     "Unique QR tied\nto booking"),
    ("🚗", "Arrives at Gate",  "Opens app,\nshows QR code"),
    ("📷", "Guard Scans",      "Any phone/tablet\nwith camera"),
    ("🟢", "Entry Granted",    "Slot marked\nOccupied instantly"),
]
for i, (icon, title, desc) in enumerate(flow_steps):
    x = 0.4 + i * 2.55
    add_rect(sl, x, 1.55, 2.3, 3.5, fill=UAE_CARD, line=UAE_GREEN, line_width=Pt(1))
    add_text(sl, icon,  x+0.85, 1.7,  0.7,  0.7, size=28, align=PP_ALIGN.CENTER)
    add_text(sl, title, x+0.1,  2.55, 2.1,  0.45, size=14, bold=True,
             color=UAE_GREEN, align=PP_ALIGN.CENTER)
    add_text(sl, desc,  x+0.1,  3.1,  2.1,  0.9, size=12,
             color=GREY_TEXT, align=PP_ALIGN.CENTER)
    if i < 4:
        add_text(sl, "→", x+2.35, 2.4, 0.25, 0.5, size=20, color=UAE_GREEN)

# Benefits row
benefits = [
    "No hardware barriers needed",
    "Works on any smartphone",
    "Complete audit trail",
    "Auto-releases after end time",
]
for i, b in enumerate(benefits):
    x = 0.5 + i * 3.22
    add_rect(sl, x, 5.3, 3.0, 0.7, fill=UAE_GREEN)
    add_text(sl, b, x+0.15, 5.38, 2.7, 0.55, size=12, bold=True,
             color=UAE_BLACK, align=PP_ALIGN.CENTER)

slide_number(sl, 6)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Tech Stack
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
dark_bg(sl)
uae_flag_bar(sl)

add_text(sl, "Technology Stack", 0.5, 0.25, 12, 0.7, size=32, bold=True, color=UAE_WHITE)
add_text(sl, "Production-grade · Cloud-native · Scalable", 0.5, 0.95, 12, 0.35,
         size=16, color=UAE_GREEN)

tech = [
    ("Flutter",          "One codebase → Web, iOS, Android",          UAE_GREEN),
    ("Firebase Auth",    "Google Sign-In · Enterprise security",       UAE_GREEN),
    ("Cloud Firestore",  "Real-time NoSQL · Auto-scaling",             UAE_GREEN),
    ("Riverpod",         "State management · Reactive UI",             UAE_GREEN),
    ("Google Maps",      "Live map integration",                       UAE_GREEN),
    ("GoRouter",         "Deep-linking · URL navigation",              UAE_GREEN),
    ("Vercel",           "Global CDN · Zero-config deployment",        UAE_RED),
    ("GitHub",           "Version control · CI/CD ready",              UAE_WHITE),
]
for i, (name, desc, color) in enumerate(tech):
    col = i % 2
    row = i // 2
    x = 0.5 + col * 6.5
    y = 1.5 + row * 1.3
    add_rect(sl, x, y, 6.1, 1.1, fill=UAE_CARD, line=color, line_width=Pt(1))
    add_rect(sl, x, y, 0.18, 1.1, fill=color)
    add_text(sl, name, x+0.35, y+0.08, 3.0, 0.45, size=15, bold=True, color=UAE_WHITE)
    add_text(sl, desc, x+0.35, y+0.55, 5.5, 0.45, size=12, color=GREY_TEXT)

slide_number(sl, 7)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — POC Results
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
dark_bg(sl)
uae_flag_bar(sl)

add_text(sl, "What This POC Proves", 0.5, 0.25, 12, 0.7, size=32, bold=True, color=UAE_WHITE)
add_text(sl, "All core capabilities validated end-to-end", 0.5, 0.95, 12, 0.35,
         size=16, color=UAE_GREEN)

results = [
    ("Real-time slot availability",       "✅ Live"),
    ("No double-booking (transactions)",  "✅ Live"),
    ("Time-based booking with expiry",    "✅ Live"),
    ("Auto-release expired bookings",     "✅ Live"),
    ("QR entry validation",               "✅ Live"),
    ("Prepaid payment flow",              "✅ Mock ready"),
    ("Multi-location support",            "✅ Live"),
    ("Web + Mobile (one codebase)",       "✅ Flutter"),
    ("Google Sign-In Auth",               "✅ Live"),
    ("Cloud Firestore backend",           "✅ Live"),
]
for i, (cap, status) in enumerate(results):
    col = i % 2
    row = i // 2
    x = 0.5 + col * 6.5
    y = 1.5 + row * 1.05
    add_rect(sl, x, y, 6.1, 0.85, fill=UAE_CARD, line=UAE_GREEN, line_width=Pt(0.75))
    add_text(sl, cap,    x+0.2,  y+0.18, 4.5, 0.5, size=13, color=UAE_WHITE)
    add_text(sl, status, x+4.9,  y+0.18, 1.1, 0.5, size=13, bold=True,
             color=UAE_GREEN, align=PP_ALIGN.RIGHT)

slide_number(sl, 8)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Business Impact
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
dark_bg(sl)
uae_flag_bar(sl)

add_text(sl, "Business Impact", 0.5, 0.25, 12, 0.7, size=32, bold=True, color=UAE_WHITE)
add_text(sl, "Operator & driver benefits", 0.5, 0.95, 12, 0.35, size=16, color=UAE_GREEN)

metrics = [
    ("0",     "Hardware needed\nfor MVP",     UAE_GREEN),
    ("100%",  "Digital audit\ntrail",         UAE_GREEN),
    ("AED 0", "Infrastructure\ncost (POC)",   UAE_RED),
    ("∞",     "Locations\nscalable",          UAE_WHITE),
]
for i, (val, label, color) in enumerate(metrics):
    x = 0.5 + i * 3.2
    add_rect(sl, x, 1.5, 3.0, 2.5, fill=UAE_CARD, line=color, line_width=Pt(1.5))
    add_text(sl, val,   x+0.2, 1.7,  2.6, 1.0, size=38, bold=True,
             color=color, align=PP_ALIGN.CENTER)
    add_text(sl, label, x+0.2, 2.75, 2.6, 0.9, size=13,
             color=GREY_TEXT, align=PP_ALIGN.CENTER)

imp = [
    ("💰 Revenue",    "Prepaid model eliminates revenue leakage. Every minute of parking captured."),
    ("👮 Operations", "Guards need zero training. Scan QR on any phone — no hardware required."),
    ("🚗 Drivers",    "Book from home, guaranteed spot, no cash, no queues at entry."),
    ("📊 Operators",  "Full booking history, occupancy rates, peak hours — all in real-time."),
]
for i, (title, desc) in enumerate(imp):
    y = 4.25 + i * 0.75
    add_rect(sl, 0.5, y, 12.3, 0.65, fill=UAE_CARD, line=UAE_GREEN, line_width=Pt(0.5))
    add_text(sl, title, 0.7,  y+0.1, 2.2, 0.45, size=13, bold=True, color=UAE_GREEN)
    add_text(sl, desc,  2.9,  y+0.1, 9.8, 0.45, size=12, color=UAE_WHITE)

slide_number(sl, 9)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Roadmap
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
dark_bg(sl)
uae_flag_bar(sl)

add_text(sl, "Roadmap to Production", 0.5, 0.25, 12, 0.7, size=32, bold=True, color=UAE_WHITE)
add_text(sl, "From POC to live product in 8-12 weeks", 0.5, 0.95, 12, 0.35,
         size=16, color=UAE_GREEN)

phases = [
    ("Phase 1\nWeeks 1-2",  UAE_RED,   [
        "Real payment gateway (Telr/Stripe)",
        "Push notifications",
        "Arabic language support",
    ]),
    ("Phase 2\nWeeks 3-5",  UAE_GREEN, [
        "Admin dashboard & reports",
        "iOS + Android apps (App Store)",
        "Multiple rate tiers per location",
    ]),
    ("Phase 3\nWeeks 6-8",  UAE_WHITE, [
        "Barrier/gate hardware integration",
        "License plate recognition (ANPR)",
        "Monthly subscriber plans",
    ]),
    ("Phase 4\nWeeks 9-12", GREY_TEXT, [
        "Multi-city rollout",
        "Fleet / corporate accounts",
        "Analytics & revenue dashboard",
    ]),
]
for i, (phase, color, items) in enumerate(phases):
    x = 0.4 + i * 3.25
    add_rect(sl, x, 1.5, 3.0, 5.5, fill=UAE_CARD, line=color, line_width=Pt(1.5))
    add_rect(sl, x, 1.5, 3.0, 0.85, fill=color)
    tc = UAE_BLACK if color == UAE_WHITE else UAE_BLACK if color == UAE_GREEN else UAE_WHITE
    add_text(sl, phase, x+0.15, 1.55, 2.7, 0.75, size=13, bold=True,
             color=tc, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        add_text(sl, f"→  {item}", x+0.2, 2.55+j*1.35, 2.6, 1.2,
                 size=12, color=UAE_WHITE)

slide_number(sl, 10)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — CTA / Close
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
dark_bg(sl)

# Full green bottom half
add_rect(sl, 0, 4.2, 13.33, 3.3, fill=UAE_GREEN)

# Red left bar
add_rect(sl, 0, 0, 0.35, 7.5, fill=UAE_RED)

# Flag top right
add_rect(sl, 11.5, 0.3, 0.1, 1.1, fill=UAE_RED)
add_rect(sl, 11.6, 0.3,  1.6, 0.37, fill=UAE_GREEN)
add_rect(sl, 11.6, 0.67, 1.6, 0.37, fill=UAE_WHITE)
add_rect(sl, 11.6, 1.04, 1.6, 0.37, fill=UAE_BLACK)

add_text(sl, "Ready to Go Live", 0.7, 1.0, 11.5, 1.2, size=46, bold=True, color=UAE_WHITE)
add_text(sl, "The foundation is built. The POC is live.", 0.7, 2.2, 11, 0.6,
         size=20, color=GREY_TEXT)
add_text(sl, "We need: a pilot location · payment gateway · your green light.",
         0.7, 2.85, 11, 0.6, size=18, color=UAE_WHITE)

add_text(sl, "🌐  Live Demo", 1.0, 4.45, 5.0, 0.55, size=16, bold=True, color=UAE_BLACK)
add_text(sl, "smartpark-lac.vercel.app", 1.0, 4.95, 5.5, 0.5, size=14, color=UAE_BLACK)

add_text(sl, "💻  Source Code", 7.0, 4.45, 5.5, 0.55, size=16, bold=True, color=UAE_BLACK)
add_text(sl, "github.com/vansan/smartpark", 7.0, 4.95, 6.0, 0.5, size=14, color=UAE_BLACK)

add_text(sl, "Questions?  Let's discuss next steps. 🚀",
         0.7, 6.1, 12.0, 0.7, size=20, bold=True, color=UAE_BLACK, align=PP_ALIGN.CENTER)

slide_number(sl, 11)

# ── Save ──────────────────────────────────────────────────────────────────────
out = r"C:\work\SmartPark\SmartPark_Stakeholder_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
